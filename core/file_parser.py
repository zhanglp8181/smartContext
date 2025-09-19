import pandas as pd
import io
import subprocess
import tempfile
import os
from docx import Document
from pptx import Presentation
from langchain_community.document_loaders import TextLoader, CSVLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from core.models import FileType
from utils.token_counter import count_tokens
import logging
import re

logger = logging.getLogger(__name__)

class FileParser:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        # 创建输出目录
        # 修复路径问题，使用项目根目录下的 output_dir
        self.output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output_dir")
        os.makedirs(self.output_dir, exist_ok=True)

    def parse_file(self, file_content, file_name: str, file_type: FileType) -> str:
        """根据文件类型解析文件内容并转换为Markdown格式"""
        try:
            logger.info(f"Parsing file {file_name} of type {file_type}")
            markdown_content = ""
            
            if file_type == FileType.TEXT or file_type == FileType.MD:
                markdown_content = self._parse_text(file_content, file_type)
            elif file_type == FileType.CSV:
                markdown_content = self._parse_csv(file_content)
            elif file_type == FileType.DOCX:
                markdown_content = self._parse_docx_with_pandoc(file_content)
            elif file_type == FileType.PPTX:
                markdown_content = self._parse_pptx_with_pandoc(file_content)
            elif file_type == FileType.XLSX:
                markdown_content = self._parse_xlsx_with_pandoc(file_content)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            # 保存Markdown文件
            self._save_markdown_file(markdown_content, file_name)
            
            return markdown_content
        except Exception as e:
            logger.error(f"Error parsing file: {str(e)}", exc_info=True)
            raise Exception(f"Error parsing file: {str(e)}")

    def _save_markdown_file(self, content: str, original_filename: str):
        """保存Markdown文件到输出目录"""
        # 生成新的文件名（保留原文件名，扩展名为.md）
        base_name = os.path.splitext(original_filename)[0]
        md_filename = f"{base_name}.md"
        md_path = os.path.join(self.output_dir, md_filename)
        
        try:
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(content)
            logger.info(f"Markdown file saved to: {md_path}")
        except Exception as e:
            logger.error(f"Failed to save markdown file: {str(e)}")

    def _parse_text(self, content: str, file_type: FileType) -> str:
        """解析文本文件并转换为Markdown格式"""
        logger.debug("Parsing text file")
        if file_type == FileType.MD:
            # 直接返回 Markdown 内容
            return content
        # 将纯文本转换为Markdown格式（添加代码块标记）
        return f"```text\n{content}\n```"

    def _parse_csv(self, content: str) -> str:
        """解析CSV文件并转换为Markdown表格"""
        logger.debug("Parsing CSV file")
        # 使用 Pandoc 将 CSV 转换为 Markdown 表格
        try:
            # 创建临时文件
            with tempfile.NamedTemporaryFile(mode='w', suffix='.csv', delete=False) as temp_csv:
                temp_csv.write(content)
                temp_csv_path = temp_csv.name
            
            # 使用 Pandoc 转换
            result = subprocess.run([
                'pandoc', '-f', 'csv', '-t', 'markdown',
                '--standalone', temp_csv_path
            ], capture_output=True, text=True, timeout=30)
            
            # 清理临时文件
            os.unlink(temp_csv_path)
            
            if result.returncode == 0:
                return result.stdout
            else:
                logger.warning(f"Pandoc CSV conversion failed: {result.stderr}")
                # 回退到原始方法
                df = pd.read_csv(io.StringIO(content))
                return df.to_markdown(index=False)
                
        except Exception as e:
            logger.warning(f"Pandoc CSV conversion failed: {str(e)}")
            # 回退到原始方法
            df = pd.read_csv(io.StringIO(content))
            return df.to_markdown(index=False)

    def _parse_docx_with_pandoc(self, content: bytes) -> str:
        """使用 Pandoc 解析 DOCX 文件为 Markdown"""
        logger.debug("Parsing DOCX file with Pandoc")
        try:
            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_docx:
                temp_docx.write(content)
                temp_docx_path = temp_docx.name
            
            # 使用 Pandoc 转换
            result = subprocess.run([
                'pandoc', '-f', 'docx', '-t', 'markdown',
                '--standalone', '--wrap=none', temp_docx_path
            ], capture_output=True, text=True, timeout=30)
            
            # 清理临时文件
            os.unlink(temp_docx_path)
            
            if result.returncode == 0:
                return result.stdout
            else:
                logger.warning(f"Pandoc DOCX conversion failed: {result.stderr}")
                # 回退到原始方法
                return self._parse_docx_fallback(content)
                
        except Exception as e:
            logger.warning(f"Pandoc DOCX conversion failed: {str(e)}")
            # 回退到原始方法
            return self._parse_docx_fallback(content)
    
    def _parse_docx_fallback(self, content: bytes) -> str:
        """DOCX 解析回退方法 - 完整转换为Markdown"""
        doc = Document(io.BytesIO(content))
        markdown_content = []
        
        # 处理段落
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # 检查段落样式
                style = paragraph.style.name
                if style.startswith('Heading'):
                    level = int(style.replace('Heading', ''))
                    markdown_content.append(f"{'#' * level} {paragraph.text}")
                else:
                    markdown_content.append(paragraph.text)
        
        # 处理表格
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                # 将表格转换为 Markdown 格式
                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                markdown_content.append(df.to_markdown(index=False))
        
        result = "\n\n".join(markdown_content)
        logger.debug(f"DOCX parsed with fallback, {len(markdown_content)} elements, {len(result)} characters")
        return result


    def _parse_pptx_with_pandoc(self, content: bytes) -> str:
        """使用 Pandoc 解析 PPTX 文件为 Markdown"""
        logger.debug("Parsing PPTX file with Pandoc")
        try:
            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
                temp_pptx.write(content)
                temp_pptx_path = temp_pptx.name
            
            # 使用 Pandoc 转换
            result = subprocess.run([
                'pandoc', '-f', 'pptx', '-t', 'markdown',
                '--standalone', '--wrap=none', temp_pptx_path
            ], capture_output=True, text=True, timeout=30)
            
            # 清理临时文件
            os.unlink(temp_pptx_path)
            
            if result.returncode == 0:
                return result.stdout
            else:
                logger.warning(f"Pandoc PPTX conversion failed: {result.stderr}")
                # 回退到原始方法
                return self._parse_pptx_fallback(content)
                
        except Exception as e:
            logger.warning(f"Pandoc PPTX conversion failed: {str(e)}")
            # 回退到原始方法
            return self._parse_pptx_fallback(content)
    
    def _parse_pptx_fallback(self, content: bytes) -> str:
        """PPTX 解析回退方法"""
        prs = Presentation(io.BytesIO(content))
        text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text.append(f"## Slide {prs.slides.index(slide) + 1}\n" + "\n".join(slide_text))
        
        result = "\n\n".join(text)
        logger.debug(f"PPTX parsed with fallback, {len(prs.slides)} slides, {len(result)} characters")
        return result

    def _parse_xlsx_with_pandoc(self, content: bytes) -> str:
        """使用 Pandoc 解析 XLSX 文件为 Markdown"""
        logger.debug("Parsing XLSX file with Pandoc")
        try:
            # 创建临时文件
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_xlsx:
                temp_xlsx.write(content)
                temp_xlsx_path = temp_xlsx.name
            
            # 使用 Pandoc 转换
            result = subprocess.run([
                'pandoc', '-f', 'xlsx', '-t', 'markdown',
                '--standalone', '--wrap=none', temp_xlsx_path
            ], capture_output=True, text=True, timeout=30)
            
            # 清理临时文件
            os.unlink(temp_xlsx_path)
            
            if result.returncode == 0:
                return result.stdout
            else:
                logger.warning(f"Pandoc XLSX conversion failed: {result.stderr}")
                # 回退到原始方法
                return self._parse_xlsx_fallback(content)
                
        except Exception as e:
            logger.warning(f"Pandoc XLSX conversion failed: {str(e)}")
            # 回退到原始方法
            return self._parse_xlsx_fallback(content)
    
    def _parse_xlsx_fallback(self, content: bytes) -> str:
        """XLSX 解析回退方法"""
        xls = pd.ExcelFile(io.BytesIO(content))
        all_text = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            all_text.append(f"## {sheet_name}\n")
            all_text.append(df.to_markdown(index=False))
        
        result = "\n\n".join(all_text)
        logger.debug(f"XLSX parsed with fallback, {len(xls.sheet_names)} sheets, {len(result)} characters")
        return result